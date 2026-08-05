"""把 PDF、Word、Excel 檔案裡的名單讀成公司紀錄。

不少公協會沒有把會員名冊做成網頁，而是掛一個 PDF 或 Excel 在網站上——常常
還要先點進某個子頁面才看得到那個連結。對這類名錄，HTML 解析完全使不上力。

## 兩種排版，兩條路

**表格**（Excel、以及 PDF 裡真的畫了格線的表）走匯入那條路：第一列是標題，
底下每一列一家公司。標題的對應（「公司名稱」「廠商名稱」「Company」都算
company_name）與「對應不到的欄位原樣保留」在 :mod:`exporter.importer` 已經
做過一次，這裡直接沿用，不重寫第二份。

**純文字**（Word、沒有格線的 PDF）走標籤那條路：「負責人︰王大明」這種寫法
由 :mod:`crawler.labels` 解析，跟舊式公會網頁用的是同一套邏輯。

## 這裡不做的事

不做 OCR。掃描成圖片的 PDF 讀不出任何文字，這時會明確回報「這個檔案沒有可
擷取的文字」，而不是安靜地回傳空清單讓使用者以為名冊是空的。
"""

from __future__ import annotations

import io
import re
from dataclasses import dataclass, field
from urllib.parse import unquote, urlsplit

from core.constants import LogCategory
from core.errors import CrawlError
from core.logging_setup import get_logger
from core.schemas import RawCompany
from crawler.labels import MIN_PAIRS, parse_record, split_cjk_english

log = get_logger(LogCategory.CRAWL)


@dataclass(frozen=True, slots=True)
class DocumentKind:
    """一種支援的檔案格式。"""

    key: str
    label: str
    suffixes: tuple[str, ...]


#: 使用者可以勾選要不要抓的檔案格式。**預設全部不勾**——下載並解析別人的
#: 檔案跟讀網頁不是同一件事，該由使用者明確決定。
DOCUMENT_KINDS: tuple[DocumentKind, ...] = (
    DocumentKind("pdf", "PDF 檔（.pdf）", (".pdf",)),
    DocumentKind("excel", "Excel 檔（.xlsx／.xls／.csv）", (".xlsx", ".xlsm", ".xls", ".csv")),
    DocumentKind("word", "Word 檔（.docx）", (".docx",)),
    DocumentKind("powerpoint", "PowerPoint 檔（.pptx）", (".pptx",)),
)

KIND_BY_KEY: dict[str, DocumentKind] = {kind.key: kind for kind in DOCUMENT_KINDS}

#: 單一檔案的大小上限。名冊再大也不會到這個數字，超過的多半是掃描檔或年報，
#: 讀進記憶體只會讓程式卡住。
MAX_DOCUMENT_BYTES = 40 * 1024 * 1024

#: 純文字模式下，一段要有這麼多字才可能是一筆紀錄。
_MIN_BLOCK_CHARS = 6


@dataclass
class DocumentResult:
    """一個檔案解析出來的結果。"""

    filename: str
    kind: str
    records: list[RawCompany] = field(default_factory=list)
    notes: list[str] = field(default_factory=list)

    @property
    def ok(self) -> bool:
        return bool(self.records)


def kind_for(url_or_name: str) -> str | None:
    """這個網址或檔名屬於哪一種支援的檔案格式；都不是就回 ``None``。"""
    path = unquote(urlsplit(url_or_name).path or url_or_name).lower()
    for kind in DOCUMENT_KINDS:
        if path.endswith(kind.suffixes):
            return kind.key
    return None


def is_wanted(url_or_name: str, wanted_kinds) -> bool:
    """使用者有沒有勾選這個檔案的格式。沒勾就完全不會被下載。"""
    kind = kind_for(url_or_name)
    return bool(kind and kind in set(wanted_kinds or ()))


# ---------------------------------------------------------------- 表格那條路


def _records_from_table(rows: list[list[str]], source_label: str) -> list[RawCompany]:
    """第一列當標題、其餘每列一家公司。認不出標題就回空清單。"""
    import pandas as pd

    from exporter.importer import rows_to_records

    if len(rows) < 2:
        return []

    header = [str(cell or "").strip() for cell in rows[0]]
    if not any(header):
        return []

    width = len(header)
    body = [
        (row + [""] * width)[:width]
        for row in rows[1:]
        if any(str(cell or "").strip() for cell in row)
    ]
    if not body:
        return []

    frame = pd.DataFrame(body, columns=_unique(header), dtype=object)
    try:
        records, _unmapped = rows_to_records(frame, source_label)
    except Exception as exc:
        # 「這張表沒有公司名稱欄」是常態（頁首、目錄、統計表都長這樣），
        # 不是錯誤——回空清單讓呼叫端去試下一張表。
        log.debug("表格無法對應成公司紀錄：{}", exc)
        return []
    return records


def _unique(header: list[str]) -> list[str]:
    """欄名去重。名冊常常有兩欄都叫「電話」，pandas 不接受重複欄名。"""
    seen: dict[str, int] = {}
    result: list[str] = []
    for index, name in enumerate(header):
        clean = name or f"欄{index + 1}"
        seen[clean] = seen.get(clean, 0) + 1
        result.append(clean if seen[clean] == 1 else f"{clean} {seen[clean]}")
    return result


# ---------------------------------------------------------------- 文字那條路


#: 兩個以上換行＝一段。名冊的每一家公司之間通常空一行。
_BLOCK_SPLIT = re.compile(r"\n\s*\n+")

#: 當成「一筆的開頭」的公司名稱行，長度上限。整段內文裡偶爾會提到公司名，
#: 那種句子通常長得多。
_NAME_LINE_MAX_CHARS = 60


def _split_blocks(text: str) -> list[str]:
    """把整份文字切成「一筆一段」。

    優先用「這一行像公司名稱」當每一筆的開頭，空行只是備援。理由是 PDF：
    ``extract_text()`` 交出來的是一行一行的文字，版面上的空白距離不會變成空行，
    所以照空行切的話整份檔案會變成一大段，只讀得到第一家公司。

    名冊的排版幾乎一定是「公司名稱一行，底下接它的聯絡資料」，所以公司名稱
    那一行就是最可靠的分隔點。
    """
    from crawler.discover import _has_company_marker

    lines = text.splitlines()
    starts = [
        index
        for index, line in enumerate(lines)
        if (stripped := line.strip())
        and len(stripped) <= _NAME_LINE_MAX_CHARS
        and _has_company_marker(stripped)
    ]
    if len(starts) >= 2:
        bounds = starts + [len(lines)]
        return ["\n".join(lines[a:b]) for a, b in zip(bounds, bounds[1:])]
    return _BLOCK_SPLIT.split(text)


def _records_from_text(text: str, source_label: str) -> list[RawCompany]:
    """把「一段一家公司」的純文字讀成紀錄。"""
    records: list[RawCompany] = []
    for block in _split_blocks(text):
        block = block.strip()
        if len(block) < _MIN_BLOCK_CHARS:
            continue

        parsed = parse_record(block.replace("\n", " "))
        if parsed.pair_count < MIN_PAIRS:
            continue

        name = parsed.fields.get("company_name", "")
        english = ""
        if not name:
            # 沒有「公司名稱︰」這種標籤時，第一個標籤之前的文字就是名稱。
            name, english = split_cjk_english(parsed.heading)
        if not name:
            continue

        values = {k: v for k, v in parsed.fields.items() if k != "company_name"}
        values.setdefault("english_name", english or None)
        records.append(
            RawCompany(
                company_name=name,
                source=source_label,
                extra_fields=dict(parsed.extra),
                **{k: v for k, v in values.items() if v},
            )
        )
    return records


# ------------------------------------------------------------------ 各格式


def _read_pdf(data: bytes, source_label: str, result: DocumentResult) -> None:
    try:
        import pdfplumber
    except ImportError as exc:  # pragma: no cover - 相依套件缺失
        raise CrawlError(
            "讀取 PDF 需要 pdfplumber 套件；請重新執行安裝檔。"
        ) from exc

    texts: list[str] = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for page in pdf.pages:
            for table in page.extract_tables() or []:
                result.records.extend(_records_from_table(table, source_label))
            texts.append(page.extract_text() or "")

    joined = "\n\n".join(texts).strip()
    if not result.records:
        result.records.extend(_records_from_text(joined, source_label))
    if not joined and not result.records:
        result.notes.append(
            "這個 PDF 沒有可擷取的文字，多半是掃描成圖片的檔案。"
            "程式不會辨識圖片裡的文字。"
        )


def _read_word(data: bytes, source_label: str, result: DocumentResult) -> None:
    try:
        import docx
    except ImportError as exc:  # pragma: no cover - 相依套件缺失
        raise CrawlError(
            "讀取 Word 需要 python-docx 套件；請重新執行安裝檔。"
        ) from exc

    document = docx.Document(io.BytesIO(data))
    for table in document.tables:
        rows = [[cell.text for cell in row.cells] for row in table.rows]
        result.records.extend(_records_from_table(rows, source_label))

    if not result.records:
        # 段落之間只接一個換行。接兩個的話每一個段落都會變成獨立的一段，
        # 而名冊的排版是「公司名稱一段、聯絡資料一段」——拆開之後名稱那一段
        # 沒有欄位、資料那一段沒有名稱，兩邊都會被丟掉。
        # 真正的分隔是檔案裡的空白段落，它本身就會補上第二個換行。
        text = "\n".join(p.text for p in document.paragraphs)
        result.records.extend(_records_from_text(text, source_label))


def _read_powerpoint(data: bytes, source_label: str, result: DocumentResult) -> None:
    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover - 相依套件缺失
        raise CrawlError(
            "讀取 PowerPoint 需要 python-pptx 套件；請重新執行安裝檔。"
        ) from exc

    presentation = Presentation(io.BytesIO(data))
    texts: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                result.records.extend(_records_from_table(rows, source_label))
            if getattr(shape, "has_text_frame", False):
                texts.append(shape.text_frame.text)
        # 一張投影片就是一個段落邊界——名冊型的簡報常常一張放一家公司。
        texts.append("")

    if not result.records:
        result.records.extend(_records_from_text("\n".join(texts), source_label))


def _read_spreadsheet(
    data: bytes, filename: str, source_label: str, result: DocumentResult
) -> None:
    import pandas as pd

    from exporter.importer import rows_to_records

    if filename.lower().endswith(".csv"):
        frames = [_read_csv(data)]
    else:
        book = pd.read_excel(io.BytesIO(data), sheet_name=None, dtype=str)
        frames = list(book.values())

    for frame in frames:
        if frame is None or frame.empty:
            continue
        try:
            records, _unmapped = rows_to_records(frame.astype(object), source_label)
        except Exception as exc:
            log.debug("工作表無法對應成公司紀錄：{}", exc)
            continue
        result.records.extend(records)


def _read_csv(data: bytes):
    import pandas as pd

    # 台灣的 CSV 很常是 Big5／CP950。UTF-8 讀不出來時逐個試，而不是直接失敗。
    for encoding in ("utf-8-sig", "utf-8", "cp950", "big5hkscs"):
        try:
            return pd.read_csv(io.BytesIO(data), dtype=str, encoding=encoding)
        except UnicodeDecodeError:
            continue
    raise CrawlError("這個 CSV 的編碼無法判讀，請另存為 UTF-8 再試一次。")


_READERS = {
    "pdf": lambda data, name, label, result: _read_pdf(data, label, result),
    "word": lambda data, name, label, result: _read_word(data, label, result),
    "powerpoint": lambda data, name, label, result: _read_powerpoint(data, label, result),
    "excel": _read_spreadsheet,
}


def extract_records(
    data: bytes, filename: str, source_label: str = "document"
) -> DocumentResult:
    """把一個檔案的位元組讀成公司紀錄。

    ``filename`` 只用來判斷格式，可以直接傳網址。格式不支援時丟
    :class:`~core.errors.CrawlError`——呼叫端本來就該先用 :func:`is_wanted`
    問過，走到這裡還不支援代表有 bug，不該安靜地當成空檔案。
    """
    kind = kind_for(filename)
    if kind is None:
        raise CrawlError(f"不支援的檔案格式：{filename}")
    if len(data) > MAX_DOCUMENT_BYTES:
        raise CrawlError(
            f"檔案超過 {MAX_DOCUMENT_BYTES // (1024 * 1024)} MB，略過不讀。"
        )

    result = DocumentResult(filename=filename, kind=kind)
    reader = _READERS[kind]
    try:
        if kind == "excel":
            reader(data, filename, source_label, result)
        else:
            reader(data, filename, source_label, result)
    except CrawlError:
        raise
    except Exception as exc:
        raise CrawlError(f"讀取 {filename} 失敗：{exc}") from exc

    # 同一份檔案裡同一家公司可能出現兩次（表格與內文各一次），先去掉重複，
    # 免得預覽的筆數看起來比實際多一倍。
    result.records = _dedupe(result.records)
    _drop_if_not_a_directory(result)
    log.info("{} 讀出 {} 筆", filename, len(result.records))
    return result


#: 讀出來的東西裡至少要有這個比例像公司名稱，才當成名冊。
#:
#: PDF 與 Word 檔裡什麼都有——章程、會議記錄、財務報表、活動照片說明。
#: 把每一段文字都當成一家公司的話，一份年報會產生幾百筆垃圾，而且混進資料庫
#: 之後很難分辨哪些是真的。寧可整份不要，也不要收一半垃圾。
_MIN_COMPANY_NAME_RATIO = 0.4

#: 少於這個筆數時不做比例判斷——三筆裡有一筆不像，比例就掉到 0.33，
#: 用比例去判斷小樣本只會誤殺。
_RATIO_SAMPLE_FLOOR = 4


def _drop_if_not_a_directory(result: DocumentResult) -> None:
    """讀出來的不像名冊就整份不要，並說明原因。"""
    from crawler.discover import _has_company_marker

    if len(result.records) < _RATIO_SAMPLE_FLOOR:
        return

    names = [record.company_name for record in result.records]
    ratio = sum(1 for name in names if _has_company_marker(name)) / len(names)
    if ratio >= _MIN_COMPANY_NAME_RATIO:
        return

    log.info(
        "{} 讀出的 {} 段文字裡只有 {:.0%} 像公司名稱，整份略過",
        result.filename, len(names), ratio,
    )
    result.records = []
    result.notes.append(
        f"這個檔案讀出來的內容不像廠商名冊（{ratio:.0%} 像公司名稱），已略過。"
        "章程、會議記錄、年報這類文件會落在這一類。"
    )


def _dedupe(records: list[RawCompany]) -> list[RawCompany]:
    seen: set[str] = set()
    unique: list[RawCompany] = []
    for record in records:
        key = record.company_name.strip()
        if not key or key in seen:
            continue
        seen.add(key)
        unique.append(record)
    return unique
