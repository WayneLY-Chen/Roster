"""PDF／Word／Excel 檔案裡的名單（``crawler.documents``）。

不少公協會沒有把會員名冊做成網頁，而是掛一個 PDF 或 Excel 在網站上，常常還
要先點進某個子頁面才看得到那個連結。
"""

from __future__ import annotations

import io

import pandas as pd
import pytest

from core.errors import CrawlError
from crawler.documents import (
    DOCUMENT_KINDS,
    MAX_DOCUMENT_BYTES,
    extract_records,
    is_wanted,
    kind_for,
)

MEMBERS = ["甲有限公司", "乙股份有限公司", "丙企業社", "丁實業有限公司", "戊工業社"]


# ------------------------------------------------------------ 格式判斷


@pytest.mark.parametrize(
    ("name", "expected"),
    [
        ("https://a.test/會員名冊.pdf", "pdf"),
        ("https://a.test/x/LIST.PDF", "pdf"),          # 副檔名大寫
        ("members.xlsx", "excel"),
        ("members.csv", "excel"),
        ("members.docx", "word"),
        ("https://a.test/list.html", None),
        ("https://a.test/", None),
    ],
)
def test_the_file_format_is_recognised_from_the_name(name, expected):
    assert kind_for(name) == expected


def test_nothing_is_downloaded_unless_the_user_ticked_that_format():
    """讀別人的檔案跟讀網頁不是同一件事，要使用者自己決定。"""
    assert is_wanted("a.pdf", ["pdf"])
    assert not is_wanted("a.pdf", ["excel"])
    assert not is_wanted("a.pdf", [])
    assert not is_wanted("a.html", ["pdf", "excel", "word"])


def test_every_offered_format_can_actually_be_read():
    """介面上勾得到的格式，程式一定要讀得懂——否則勾了也沒用。"""
    from crawler.documents import _READERS

    for kind in DOCUMENT_KINDS:
        assert kind.key in _READERS
        assert kind.label and kind.suffixes


# ------------------------------------------------------------ Excel


def _excel(rows: list[dict]) -> bytes:
    buffer = io.BytesIO()
    pd.DataFrame(rows).to_excel(buffer, index=False)
    return buffer.getvalue()


def test_an_excel_member_list_becomes_records():
    data = _excel(
        [{"公司名稱": name, "電話": "02-1111-2222", "會員代表": "王大明"} for name in MEMBERS]
    )
    result = extract_records(data, "members.xlsx", "test")

    assert [r.company_name for r in result.records] == MEMBERS
    assert result.records[0].phone == "02-1111-2222"


def test_a_column_with_no_home_is_kept_as_a_free_form_field():
    """欄位對應與「對應不到的原樣保留」跟匯入走同一套，不重寫第二份。"""
    data = _excel([{"公司名稱": name, "入會年月日": "1970"} for name in MEMBERS])
    result = extract_records(data, "members.xlsx", "test")

    assert result.records[0].extra_fields["入會年月日"] == "1970"


def test_a_csv_in_a_legacy_encoding_is_still_readable():
    """台灣的 CSV 很常是 Big5／CP950。"""
    text = "公司名稱,電話\n" + "\n".join(f"{name},02-1234" for name in MEMBERS)
    result = extract_records(text.encode("cp950"), "members.csv", "test")

    assert [r.company_name for r in result.records] == MEMBERS


# ------------------------------------------------------------ Word


def _word(paragraphs: list[str]) -> bytes:
    import docx

    document = docx.Document()
    for text in paragraphs:
        document.add_paragraph(text)
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def test_a_word_member_list_becomes_records():
    paragraphs: list[str] = []
    for index, name in enumerate(MEMBERS, start=1):
        paragraphs.append(name)
        paragraphs.append(f"負責人︰ 王大明   電話︰ 02-1111-{index:04d}   地址︰ 台北市中山區")
        paragraphs.append("")

    result = extract_records(_word(paragraphs), "members.docx", "test")

    assert [r.company_name for r in result.records] == MEMBERS
    assert result.records[0].contact_person == "王大明"
    assert result.records[0].phone == "02-1111-0001"


def test_the_company_name_and_its_details_are_not_split_apart():
    """名冊的排版是「公司名稱一段、聯絡資料一段」。拆開的話名稱那一段沒有
    欄位、資料那一段沒有名稱，兩邊都會被丟掉。"""
    result = extract_records(
        _word(["甲有限公司", "電話︰ 02-1   傳真︰ 02-2"]), "x.docx", "test"
    )
    assert len(result.records) == 1
    assert result.records[0].company_name == "甲有限公司"
    assert result.records[0].fax == "02-2"


# ------------------------------------------------------------ PDF


def _pdf(lines: list[str]) -> bytes:
    """最小可讀的 PDF。只要能讓 pdfplumber 取出這幾行文字就夠了。"""

    def escape(text: str) -> str:
        return text.replace("\\", "").replace("(", "").replace(")", "")

    content = "BT /F1 12 Tf 40 760 Td 14 TL\n" + "\n".join(
        f"({escape(line)}) Tj T*" for line in lines
    ) + "\nET"
    stream = content.encode("latin-1")
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 595 842] "
        b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length %d >>\nstream\n" % len(stream) + stream + b"\nendstream",
    ]
    out = bytearray(b"%PDF-1.4\n")
    offsets: list[int] = []
    for index, body in enumerate(objects, start=1):
        offsets.append(len(out))
        out += b"%d 0 obj\n" % index + body + b"\nendobj\n"
    xref = len(out)
    out += b"xref\n0 %d\n" % (len(objects) + 1) + b"0000000000 65535 f \n"
    for offset in offsets:
        out += b"%010d 00000 n \n" % offset
    out += b"trailer\n<< /Size %d /Root 1 0 R >>\nstartxref\n%d\n%%EOF\n" % (
        len(objects) + 1,
        xref,
    )
    return bytes(out)


#: 這個最小 PDF 只嵌得起 Helvetica，所以測試資料用英文公司名。
_EN_MEMBERS = [
    "ACME Trading Co., Ltd.",
    "Bright Star Industries Inc.",
    "Pacific Supply Co., Ltd.",
    "Formosa Machinery Corp.",
]


def test_a_pdf_member_list_becomes_records():
    lines: list[str] = []
    for index, name in enumerate(_EN_MEMBERS, start=1):
        lines.append(name)
        lines.append(f"Contact: Mr. Wang   Tel: 02-1111-{index:04d}")

    result = extract_records(_pdf(lines), "members.pdf", "test")

    assert [r.company_name for r in result.records] == _EN_MEMBERS
    assert result.records[0].phone == "02-1111-0001"


def test_a_pdf_without_blank_lines_is_still_split_per_company():
    """PDF 的 extract_text() 交出來的是一行一行的文字，版面上的空白距離不會
    變成空行——照空行切的話整份檔案會變成一大段，只讀得到第一家公司。"""
    lines: list[str] = []
    for name in _EN_MEMBERS:
        lines.append(name)
        lines.append("Contact: Mr. Wang   Tel: 02-1111-2222")

    result = extract_records(_pdf(lines), "members.pdf", "test")
    assert len(result.records) == len(_EN_MEMBERS)


def test_a_scanned_pdf_says_so_instead_of_looking_empty():
    """整份沒有文字時要講出來，不要安靜地回空清單讓人以為名冊是空的。"""
    result = extract_records(_pdf([]), "scan.pdf", "test")

    assert not result.records
    assert any("掃描" in note for note in result.notes)


# --------------------------------------------------- 不是名冊的檔案


def test_a_document_that_is_not_a_member_list_is_skipped_entirely():
    """PDF 與 Word 檔裡什麼都有——章程、會議記錄、財務報表。把每一段都當成
    一家公司的話，一份年報會產生幾百筆垃圾，而且混進資料庫之後很難分辨。"""
    data = _excel(
        [
            {"公司名稱": text, "電話": "02-1"}
            for text in [
                "第一條 本會定名為",
                "第二條 本會以聯絡同業感情為宗旨",
                "第三條 本會會址設於",
                "第四條 本會任務如下",
                "第五條 入會資格",
            ]
        ]
    )
    result = extract_records(data, "章程.xlsx", "test")

    assert result.records == []
    assert any("不像廠商名冊" in note for note in result.notes)


def test_a_short_list_is_not_judged_by_ratio():
    """三筆裡有一筆不像，比例就掉到 0.33——用比例去判斷小樣本只會誤殺。"""
    data = _excel([{"公司名稱": "甲有限公司"}, {"公司名稱": "台積電"}])
    result = extract_records(data, "small.xlsx", "test")

    assert len(result.records) == 2


def test_the_same_company_listed_twice_is_only_kept_once():
    data = _excel([{"公司名稱": name} for name in MEMBERS + MEMBERS])
    result = extract_records(data, "members.xlsx", "test")

    assert [r.company_name for r in result.records] == MEMBERS


# ------------------------------------------------------------ 防呆


def test_an_unsupported_format_is_an_error_not_an_empty_result():
    """呼叫端本來就該先問過 is_wanted；走到這裡還不支援代表有 bug。"""
    with pytest.raises(CrawlError):
        extract_records(b"x", "notes.txt", "test")


def test_an_oversized_file_is_refused_before_it_is_parsed():
    with pytest.raises(CrawlError, match="MB"):
        extract_records(b"x" * (MAX_DOCUMENT_BYTES + 1), "big.pdf", "test")


def test_a_corrupt_file_reports_which_file_failed():
    with pytest.raises(CrawlError, match="members.pdf"):
        extract_records(b"this is not a pdf", "members.pdf", "test")


# --------------------------------------------- 爬取時跟進頁面上的檔案連結


def _crawl(tmp_config, html: str, files: dict[str, bytes], kinds: list[str], **overrides):
    """爬一頁，頁面上掛著幾個檔案連結。回傳收到的紀錄。"""
    import httpx

    from core.config import PaginationRule, SourceConfig
    from crawler.fetcher import HttpxFetcher
    from crawler.robots import RobotsPolicy
    from crawler.sources.generic_html import GenericHtmlSource

    def handler(request: httpx.Request) -> httpx.Response:
        body = files.get(request.url.path)
        if body is not None:
            return httpx.Response(200, content=body)
        return httpx.Response(200, text=html, headers={"Content-Type": "text/html"})

    client = httpx.Client(transport=httpx.MockTransport(handler))
    config = SourceConfig(
        name="docs",
        type="generic_html",
        start_url="https://a.test/list",
        list_selector="div.item",
        fields={"company_name": {"selector": "h3"}},
        pagination=PaginationRule(type="none"),
        document_kinds=kinds,
        **overrides,
    )
    source = GenericHtmlSource(
        config,
        fetcher=HttpxFetcher(
            config=tmp_config, robots=RobotsPolicy("ua", enabled=False), client=client
        ),
        config=tmp_config,
    )
    return list(source.iter_pages())[0].records


_PAGE_WITH_LINK = """<html><body>
  <div class='item'><h3>網頁上的甲有限公司</h3></div>
  <div class='item'><h3>網頁上的乙有限公司</h3></div>
  <div class='item'><h3>網頁上的丙有限公司</h3></div>
  <a href="/會員名冊.xlsx">會員名冊下載</a>
</body></html>"""


def test_a_linked_file_is_not_touched_unless_the_user_ticked_it(tmp_config):
    """沒勾就一個檔案都不下載——這是預設行為。"""
    files = {"/會員名冊.xlsx": _excel([{"公司名稱": name} for name in MEMBERS])}
    records = _crawl(tmp_config, _PAGE_WITH_LINK, files, kinds=[])

    assert all(not r.company_name.startswith(("甲", "乙")) or r.company_name.startswith("網頁上的")
               for r in records)
    assert len(records) == 3


def test_a_ticked_file_is_read_and_its_records_are_collected(tmp_config):
    files = {"/會員名冊.xlsx": _excel([{"公司名稱": name} for name in MEMBERS])}
    records = _crawl(tmp_config, _PAGE_WITH_LINK, files, kinds=["excel"])

    names = [r.company_name for r in records]
    assert len(names) == 3 + len(MEMBERS)
    for name in MEMBERS:
        assert name in names


def test_the_file_records_remember_which_file_they_came_from(tmp_config):
    files = {"/會員名冊.xlsx": _excel([{"公司名稱": name} for name in MEMBERS])}
    records = _crawl(tmp_config, _PAGE_WITH_LINK, files, kinds=["excel"])

    from_file = [r for r in records if r.company_name in MEMBERS]
    assert all(r.source_url.endswith(".xlsx") for r in from_file)


def test_the_number_of_files_read_has_a_ceiling(tmp_config):
    """每一個檔案都是一次完整下載，要有上限。"""
    html = "<html><body><div class='item'><h3>網頁上的公司</h3></div>" + "".join(
        f"<a href='/f{i}.xlsx'>名冊 {i}</a>" for i in range(10)
    ) + "</body></html>"
    # 每個檔案裝不同的公司，才看得出到底讀了幾個檔案。
    files = {
        f"/f{i}.xlsx": _excel([{"公司名稱": f"第{i}檔的{name}"} for name in MEMBERS])
        for i in range(10)
    }

    records = _crawl(tmp_config, html, files, kinds=["excel"], max_documents=2)

    from_files = [r for r in records if r.company_name.startswith("第")]
    assert len(from_files) == 2 * len(MEMBERS)


# ------------------------------------------------------------ PowerPoint


def _pptx(slides: list[list[str]]) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    blank = presentation.slide_layouts[6]
    for lines in slides:
        slide = presentation.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
        frame = box.text_frame
        frame.text = lines[0]
        for line in lines[1:]:
            frame.add_paragraph().text = line
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def test_a_powerpoint_member_list_becomes_records():
    """公會的簡報型名冊常常一張投影片放一家公司。"""
    slides = [
        [name, f"負責人︰ 王大明   電話︰ 02-1111-{i:04d}"]
        for i, name in enumerate(MEMBERS, start=1)
    ]
    result = extract_records(_pptx(slides), "members.pptx", "test")

    assert [r.company_name for r in result.records] == MEMBERS
    assert result.records[0].phone == "02-1111-0001"


def test_powerpoint_is_offered_and_readable():
    assert kind_for("https://a.test/名冊.pptx") == "powerpoint"
    assert is_wanted("a.pptx", ["powerpoint"])
    assert not is_wanted("a.pptx", ["pdf"])
